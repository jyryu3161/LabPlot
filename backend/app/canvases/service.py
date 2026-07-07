from __future__ import annotations

import base64
import copy
import hashlib
import io
import json
import math
import os
import re
import shutil
import subprocess
import tempfile
import uuid
from xml.sax.saxutils import escape as _xml_escape

from sqlalchemy import or_
from sqlalchemy.orm import Session, joinedload

from app.canvases import images as canvas_images
from app.canvases.models import Canvas, CanvasPanel
from app.canvases.schemas import CanvasCreate, PanelCreate, PreviewRenderRequest
from app.common import storage
from app.common.encryption import decrypt_private_bytes
from app.common.exceptions import AppError, BadRequestError, NotFoundError
from app.config import settings
from app.datasets import service as ds_service
from app.datasets.models import Dataset
from app.figures import service as figures_service
from app.figures.models import Figure, FigureVersion
from app.r_engine import renderer
from app.r_engine.presets import JOURNAL_SPECS, PRESETS
from app.r_engine.templates import scale_editable_axes

# mm clamps (design §5): canvas 20-500 mm/side; panel 10-500 mm/side. Enforced
# again defensively here even though the request schema already validates the
# range, so a degenerate value can never reach the R device / a stored row.
_CANVAS_MM_MIN = 20.0
_CANVAS_MM_MAX = 500.0
_PANEL_MM_MIN = 10.0
_PANEL_MM_MAX = 500.0
# Off-sheet placement envelope for a panel's top-left (mm). A panel may be
# parked beside/off the A4 sheet and pulled back later; bounded so it can't be
# flung unrecoverably far. MUST match the frontend OFF_SHEET_MIN/MAX_MM in
# CanvasEditor.tsx so a dragged position never re-snaps server-side.
_PANEL_POS_MIN = -_CANVAS_MM_MAX  # -500
_PANEL_POS_MAX = 2 * _CANVAS_MM_MAX  # 1000

# Allowed canvas backgrounds (design §2). Anything else falls back to "white".
_CANVAS_BACKGROUNDS = {"white", "transparent"}

# Ephemeral preview cache namespace (design §4). Mirrors the existing figures
# storage layout: local under static/figures/canvases/preview, object storage
# under the "figures/canvases/preview" key prefix.
_PREVIEW_PARTS = ("figures", "canvases", "preview")

# Imported external images (SVG/PNG/JPEG panels). Stored under the figures
# root — local: static/figures/canvases/imports, object storage: the
# "figures/canvases/imports" key prefix — so the existing /static mount and the
# /api/assets figures/ allow-list serve them with no new route (public-by-key,
# same grade as rendered figures; grilling Q3). image_key is the RELATIVE part
# ("canvases/imports/<hex32>.<ext>"), backend-agnostic.
_IMAGE_KEY_RE = re.compile(r"^canvases/imports/[0-9a-f]{32}\.(png|jpg|svg)$")
_IMAGE_MEDIA_TYPES = {"png": "image/png", "jpg": "image/jpeg", "svg": "image/svg+xml"}


def _image_ref(image_key: str) -> str:
    """storage.read_bytes-able ref for an import blob (local path or s3 URI)."""
    if storage.object_storage_enabled():
        return storage.object_uri(storage.object_key("figures", image_key))
    return os.path.join(settings.figures_dir, image_key)


def _image_url(image_key: str) -> str | None:
    return figures_service._url(_image_ref(image_key))


def _image_ext(image_key: str) -> str:
    return image_key.rsplit(".", 1)[-1].lower()


def _clamp(value: float, low: float, high: float) -> float:
    return max(low, min(high, float(value)))


def _dataset_content_hash(ds: Dataset) -> str:
    """Stable content hash of the dataset (design §4).

    The Dataset model stores no precomputed content hash, so hash the decrypted
    file bytes. This keys the preview cache to the *data* — a new upload/replace
    produces a different hash and misses the stale cache.
    """
    raw = decrypt_private_bytes(storage.read_bytes(ds.file_path))
    return hashlib.sha256(raw).hexdigest()


def _local_cache_path(digest: str) -> str:
    return os.path.join(settings.figures_dir, "canvases", "preview", f"{digest}.svg")


def _local_layout_path(digest: str) -> str:
    return os.path.join(settings.figures_dir, "canvases", "preview", f"{digest}.layout.json")


def _object_cache_ref(digest: str) -> str:
    return storage.object_uri(storage.object_key(*_PREVIEW_PARTS, f"{digest}.svg"))


def _read_layout_local(digest: str) -> dict | None:
    """Parse the cached preview sidecar (series_hex / legend_keys / panel_px in
    the preview SVG's own px coords) — used by the canvas color editor for
    scoped instant recolor. Best-effort: any failure returns None."""
    try:
        path = _local_layout_path(digest)
        if os.path.exists(path):
            with open(path, encoding="utf-8") as fh:
                data = json.load(fh)
            return data if isinstance(data, dict) else None
    except Exception:
        return None
    return None


def _read_layout_object(digest: str) -> dict | None:
    try:
        ref = storage.object_uri(storage.object_key(*_PREVIEW_PARTS, f"{digest}.layout.json"))
        if storage.exists(ref):
            data = json.loads(storage.read_bytes(ref).decode("utf-8"))
            return data if isinstance(data, dict) else None
    except Exception:
        return None
    return None


def render_preview(db: Session, owner_id: uuid.UUID, req: PreviewRenderRequest) -> dict:
    """Ephemeral single-SVG preview render (design decision 4, §3, §4).

    Thin public wrapper over :func:`_render_preview_ref` that maps the readable
    SVG storage ref to a served URL. External contract unchanged
    (``{svg_url, cached, layout}``).
    """
    ref, layout, cached = _render_preview_ref(db, owner_id, req)
    return {"svg_url": figures_service._url(ref), "cached": cached, "layout": layout}


def _render_preview_ref(db: Session, owner_id: uuid.UUID, req: PreviewRenderRequest) -> tuple[str, dict | None, bool]:
    """Ephemeral single-SVG preview render → (svg_ref, layout, cached).

    Renders the figure's current (or pinned) version at a custom physical size
    with an optional color/base_size overlay, WITHOUT creating a FigureVersion
    and WITHOUT touching the rerender() path. Content-hash cached: an identical
    request returns the cached SVG with cached=True and never re-renders.

    ``svg_ref`` is a reference readable by ``storage.read_bytes`` (a local
    filesystem path or an ``s3://`` object URI) — the canvas export composer
    reuses this to obtain each panel's raw physical-size vector SVG.
    """
    # 1. Owner/project-scoped figure access (404/403 as usual).
    fig = figures_service.get_figure(db, req.figure_id, owner_id)

    # Pick the version: explicit pin (must belong to the figure) else current.
    if req.version_id is not None:
        version = figures_service.get_version(fig, req.version_id)
    elif fig.current_version_id:
        version = figures_service.get_version(fig, fig.current_version_id)
    else:
        version = figures_service._current_or_latest_version(fig)
    if version is None:
        raise BadRequestError("Figure has no version to preview", error_code="NO_VERSION")

    plot_type = fig.plot_type
    ds = ds_service.get_dataset(db, fig.dataset_id, owner_id)
    valid_columns = figures_service._dataset_column_names(ds)

    # 2. Merge the overlay (series_styles / category_colors / base_size) ON TOP
    #    of the version's stored options, then run the whole thing through the
    #    figures allow-list so ONLY allow-listed keys reach R.
    merged = dict(version.options or {})
    overlay = req.options_overlay
    if overlay is not None:
        if overlay.series_styles is not None:
            merged["series_styles"] = overlay.series_styles
        if overlay.category_colors is not None:
            merged["category_colors"] = overlay.category_colors
        if overlay.base_size is not None:
            merged["base_size"] = overlay.base_size

    options = figures_service.sanitize_options(plot_type, merged, valid_columns)

    # Custom physical size -> inches for the R device (mm/25.4). base_size stays
    # absolute pt regardless of (w,h) — the re-layout guarantee (§5).
    w_mm = _clamp(req.width_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
    h_mm = _clamp(req.height_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
    options["size"] = "custom"
    options["width_in"] = w_mm / 25.4
    options["height_in"] = h_mm / 25.4

    preset = version.style_preset if version.style_preset in PRESETS else fig.style_preset
    if preset not in PRESETS:
        preset = "nature"

    mapping = version.mapping or {}

    # 3. Content-hash cache key (§4 preview cache).
    key_material = json.dumps(
        {
            "dataset_content_hash": _dataset_content_hash(ds),
            "plot_type": plot_type,
            "mapping": mapping,
            "options": options,
            "style_preset": preset,
            "w_mm": round(w_mm, 2),
            "h_mm": round(h_mm, 2),
            # Sidecar schema version: bump when layout_export gains keys, so
            # cached previews from older renderer builds (whose .layout.json
            # lacks the new hit boxes) stop being served. v4 = facet axis unions
            # exclude zero-size interior cells (v3 boxes engulfed panel rows).
            "sidecar_v": 4,
        },
        sort_keys=True,
        default=str,
    )
    digest = hashlib.sha256(key_material.encode("utf-8")).hexdigest()

    # Request-time augmentation (NOT cached): whether the backend will actually
    # emit a tick/format/reverse scale layer per aesthetic for THIS figure —
    # the axis popover hides those controls when false. Computed here so stale
    # cached sidecars get the flags too, with the CURRENT gate logic.
    _scale_flags = scale_editable_axes(plot_type, mapping, options)

    def _aug(layout: dict | None) -> dict | None:
        if isinstance(layout, dict):
            layout = dict(layout)
            layout["scale_editable_x"] = _scale_flags["x"]
            layout["scale_editable_y"] = _scale_flags["y"]
        return layout

    # Cache hit -> return without rendering, no version, no artifact written.
    if storage.object_storage_enabled():
        cache_ref = _object_cache_ref(digest)
        if storage.exists(cache_ref):
            return cache_ref, _aug(_read_layout_object(digest)), True
    else:
        cache_path = _local_cache_path(digest)
        if os.path.exists(cache_path) and os.path.getsize(cache_path) > 0:
            return cache_path, _aug(_read_layout_local(digest)), True

    # 4. Render (sandbox intact via renderer.render). Copy ONLY the SVG to the
    #    cache path; never persist png/version/etc. On failure raise RENDER_FAILED
    #    with no partial artifact.
    df = ds_service.load_dataframe(ds)
    with tempfile.TemporaryDirectory(prefix="labplot_preview_") as tmp_out_dir:
        res = renderer.render(plot_type, mapping, options, preset, df, tmp_out_dir)
        if not res.success:
            raise BadRequestError(figures_service._friendly_error(res.log), error_code="RENDER_FAILED")
        svg_src = (res.outputs or {}).get("svg")
        if not svg_src or not os.path.exists(svg_src):
            raise BadRequestError("Preview render produced no SVG", error_code="RENDER_FAILED")

        # Parse the sidecar produced for THIS preview render (series_hex /
        # legend_keys / panel_px in the preview SVG's own px), so the color
        # editor can scope its instant recolor. Best-effort.
        layout = None
        layout_src = (res.outputs or {}).get("layout")
        if layout_src and os.path.exists(layout_src):
            try:
                with open(layout_src, encoding="utf-8") as fh:
                    parsed = json.load(fh)
                layout = parsed if isinstance(parsed, dict) else None
            except Exception:
                layout = None

        if storage.object_storage_enabled():
            key = storage.object_key(*_PREVIEW_PARTS, f"{digest}.svg")
            cache_ref = storage.upload_file(svg_src, key, content_type="image/svg+xml")
            if layout is not None and layout_src:
                try:
                    storage.upload_file(layout_src, storage.object_key(*_PREVIEW_PARTS, f"{digest}.layout.json"), content_type="application/json")
                except Exception:
                    pass
            return cache_ref, _aug(layout), False

        cache_path = _local_cache_path(digest)
        os.makedirs(os.path.dirname(cache_path), exist_ok=True)
        shutil.copyfile(svg_src, cache_path)
        if layout is not None and layout_src:
            try:
                shutil.copyfile(layout_src, _local_layout_path(digest))
            except Exception:
                pass
        return cache_path, _aug(layout), False


# ============================================================================
# M2 — canvas + panel CRUD (design §3). Plain owner/project-scoped CRUD over the
# `canvases` / `canvas_panels` tables. NONE of this touches the R engine or
# creates a FigureVersion: geometry/placement is canvas-owned (§1) and a panel
# PATCH only mutates a canvas_panels row.
# ============================================================================


# ---------------------------------------------------------------- retrieval
def get_canvas(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID, write: bool = False) -> Canvas:
    """Owner/project-scoped canvas fetch, mirroring figures.get_figure exactly.

    Visible when the caller owns it OR has project access; a write additionally
    requires project write access when the canvas is project-scoped. Anything
    else is a 404 (never leak existence across tenants).
    """
    from app.projects import service as project_service

    canvas = (
        db.query(Canvas)
        .options(joinedload(Canvas.panels))
        .filter(Canvas.id == canvas_id)
        .first()
    )
    if not canvas or (
        canvas.owner_id != owner_id
        and not project_service.can_access_project(db, canvas.project_id, owner_id)
    ):
        raise NotFoundError("Canvas", str(canvas_id))
    if write and canvas.owner_id != owner_id:
        project_service.require_project_write(db, canvas.project_id, owner_id)
    return canvas


def _get_panel(canvas: Canvas, panel_id: uuid.UUID) -> CanvasPanel:
    """A panel is only reachable through its canvas (panel.canvas_id == canvas.id)."""
    for panel in canvas.panels:
        if panel.id == panel_id:
            return panel
    raise NotFoundError("CanvasPanel", str(panel_id))


def _figure_current_versions(db: Session, figure_ids: list[uuid.UUID],
                             owner_id: uuid.UUID) -> dict[uuid.UUID, uuid.UUID | None]:
    """Map figure_id -> figure.current_version_id for follow-latest resolution.

    Filtered to figures the CALLER can access (same predicate as
    figures_service.get_figure: owner match OR accessible project): a figure
    ABSENT from the map is inaccessible and must fail closed
    (effective_version_id=None, no native size in _panel_response) — a
    caller-owned canvas (e.g. a U9 duplicate of a project canvas) would
    otherwise keep resolving live version ids/metadata of figures the caller
    can't, or can no longer, access."""
    from app.projects import service as project_service

    ids = list({fid for fid in figure_ids if fid is not None})
    if not ids:
        return {}
    accessible = project_service.accessible_project_ids(db, owner_id)
    rows = (
        db.query(Figure.id, Figure.current_version_id)
        .filter(Figure.id.in_(ids),
                or_(Figure.owner_id == owner_id, Figure.project_id.in_(accessible)))
        .all()
    )
    return {fid: cvid for fid, cvid in rows}


def _version_native_sizes(db: Session, version_ids: list[uuid.UUID | None]) -> dict[uuid.UUID, tuple[float | None, float | None]]:
    """Map version_id -> (native_width_mm, native_height_mm) from its options.

    One query for all effective versions of a canvas; feeds the panel
    serializer so the editor can offer original-size placement/reset."""
    ids = list({vid for vid in version_ids if vid is not None})
    if not ids:
        return {}
    rows = db.query(FigureVersion.id, FigureVersion.options).filter(FigureVersion.id.in_(ids)).all()
    return {vid: figures_service.native_size_mm(options) for vid, options in rows}


# ---------------------------------------------------------------- serialization
def _panel_response(panel: CanvasPanel, current_map: dict[uuid.UUID, uuid.UUID | None],
                    native_map: dict[uuid.UUID, tuple[float | None, float | None]] | None = None) -> dict:
    # effective_version_id (§3): the pin if set, else the figure's current
    # version (None if the figure has no version yet). Resolved here so the
    # editor needs no extra round-trip. render_url is the committed derived-cache
    # artifact (§4), populated by later milestones; None for pure CRUD.
    # ABSENCE from current_map means the figure is INACCESSIBLE to the caller
    # (see _figure_current_versions) — fail closed, including the pin, so no
    # version id or native-size metadata of an unreadable figure leaks.
    if panel.image_key:
        # Imported-image panel: no figure/version machinery. render_url serves
        # the stored (sanitized) blob directly; native size was computed once
        # at upload from the image's own dimensions/DPI.
        effective = None
        nw, nh = panel.image_native_width_mm, panel.image_native_height_mm
        render_url = _image_url(panel.image_key)
    elif panel.figure_id not in current_map:
        effective = None
        nw, nh = None, None
        render_url = None
    else:
        effective = panel.pinned_version_id or current_map.get(panel.figure_id)
        native = (native_map or {}).get(effective) if effective else None
        nw, nh = native if native else (None, None)
        render_url = None
    return {
        "native_width_mm": nw,
        "native_height_mm": nh,
        "id": panel.id,
        "canvas_id": panel.canvas_id,
        "figure_id": panel.figure_id,
        "image_key": panel.image_key,
        "pinned_version_id": panel.pinned_version_id,
        "x_mm": panel.x_mm,
        "y_mm": panel.y_mm,
        "width_mm": panel.width_mm,
        "height_mm": panel.height_mm,
        "z_order": panel.z_order,
        "label": panel.label,
        "label_visible": panel.label_visible,
        "created_at": panel.created_at,
        "updated_at": panel.updated_at,
        "effective_version_id": effective,
        "render_url": render_url,
    }


def _canvas_detail(db: Session, canvas: Canvas, owner_id: uuid.UUID) -> dict:
    current_map = _figure_current_versions(db, [p.figure_id for p in canvas.panels], owner_id)
    native_map = _version_native_sizes(
        db, [p.pinned_version_id or current_map.get(p.figure_id) for p in canvas.panels])
    # Ordered by z_order; ties broken by id for stable paint order (§2).
    panels = sorted(canvas.panels, key=lambda p: (p.z_order, str(p.id)))
    return {
        "id": canvas.id,
        "name": canvas.name,
        "description": canvas.description,
        "owner_id": canvas.owner_id,
        "project_id": canvas.project_id,
        "width_mm": canvas.width_mm,
        "height_mm": canvas.height_mm,
        "preset": canvas.preset,
        "background": canvas.background,
        "export_snapshot": canvas.export_snapshot,
        "created_at": canvas.created_at,
        "updated_at": canvas.updated_at,
        "panels": [_panel_response(p, current_map, native_map) for p in panels],
        "annotations": canvas.annotations or [],
        "annotations_rev": canvas.annotations_rev or 0,
    }


def _canvas_list_item(canvas: Canvas) -> dict:
    return {
        "id": canvas.id,
        "name": canvas.name,
        "project_id": canvas.project_id,
        "width_mm": canvas.width_mm,
        "height_mm": canvas.height_mm,
        "panel_count": len(canvas.panels),
        "updated_at": canvas.updated_at,
    }


# ---------------------------------------------------------------- canvases
def list_canvases(db: Session, owner_id: uuid.UUID, project_id: uuid.UUID | None = None) -> list[dict]:
    """Owner's canvases plus project-accessible ones, newest first (§3)."""
    from app.projects import service as project_service

    q = db.query(Canvas).options(joinedload(Canvas.panels))
    if project_id is not None:
        project_service.get_project_model(db, project_id, owner_id)
        q = q.filter(Canvas.project_id == project_id)
    else:
        ids = project_service.accessible_project_ids(db, owner_id)
        q = q.filter(or_(Canvas.owner_id == owner_id, Canvas.project_id.in_(ids)))
    rows = q.order_by(Canvas.updated_at.desc()).all()
    return [_canvas_list_item(c) for c in rows]


def create_canvas(db: Session, owner_id: uuid.UUID, data: CanvasCreate) -> dict:
    from app.projects import service as project_service

    if data.project_id is not None:
        project_service.require_project_write(db, data.project_id, owner_id)
    background = data.background if data.background in _CANVAS_BACKGROUNDS else "white"
    canvas = Canvas(
        owner_id=owner_id,
        project_id=data.project_id,
        name=data.name,
        description=data.description,
        width_mm=_clamp(data.width_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX),
        height_mm=_clamp(data.height_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX),
        preset=data.preset,
        background=background,
    )
    db.add(canvas)
    db.commit()
    return canvas_detail(db, canvas.id, owner_id)


def canvas_detail(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID) -> dict:
    canvas = get_canvas(db, canvas_id, owner_id)
    return _canvas_detail(db, canvas, owner_id)


def update_canvas(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID, data: dict) -> dict:
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    if data.get("name") is not None:
        canvas.name = data["name"]
    if "description" in data:
        canvas.description = data["description"]
    if data.get("width_mm") is not None:
        canvas.width_mm = _clamp(data["width_mm"], _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    if data.get("height_mm") is not None:
        canvas.height_mm = _clamp(data["height_mm"], _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    if data.get("background") is not None:
        canvas.background = data["background"] if data["background"] in _CANVAS_BACKGROUNDS else "white"
    if "preset" in data:
        canvas.preset = data["preset"]
    if "project_id" in data:
        # Owner-only (grilling Q6): attach/move/detach all change which team
        # can see the canvas — an editor could otherwise privatize it.
        if canvas.owner_id != owner_id:
            # NOTE: AppError's signature is (status_code, detail, error_code) —
            # detail must not be passed positionally first (raises TypeError
            # at raise time -> 500 instead of the intended status).
            raise AppError(status_code=403,
                           detail="Only the canvas owner can move it between projects",
                           error_code="OWNER_ONLY")
        new_project_id = data["project_id"]
        if new_project_id is not None:
            from app.projects import service as project_service
            project_service.require_project_write(db, new_project_id, owner_id)
        canvas.project_id = new_project_id
    # U8: None (absent OR explicit null) = leave unchanged; [] = clear all.
    # Any editor with write access (same permission model as panels — write=True
    # above already required it) may edit annotations.
    if data.get("annotations") is not None:
        # Optimistic-concurrency guard (mirrors figures RerenderRequest.
        # base_version_id -> 409 VERSION_CONFLICT): annotations are replaced
        # whole-array, so two concurrent editors would otherwise silently
        # last-write-wins each other's objects. When the client supplies the
        # rev it based its edit on and it no longer matches, 409 instead of
        # destroying the other editor's work. Omitted (None) keeps legacy
        # last-write-wins for direct API callers.
        base_rev = data.get("base_annotations_rev")
        if base_rev is not None and base_rev != (canvas.annotations_rev or 0):
            raise AppError(
                status_code=409,
                detail="Canvas annotations were modified by another editor; reload and retry",
                error_code="ANNOTATIONS_CONFLICT",
            )
        canvas.annotations = _sanitize_annotations(data["annotations"], canvas.width_mm, canvas.height_mm)
        canvas.annotations_rev = (canvas.annotations_rev or 0) + 1
    db.commit()
    return canvas_detail(db, canvas_id, owner_id)


def delete_canvas(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID) -> None:
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    db.delete(canvas)  # FK ON DELETE CASCADE + orm cascade removes panels
    db.commit()


def duplicate_canvas(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID) -> dict:
    """Deep-copy a canvas (+ panels + annotations) into a new canvas owned by
    the CALLER (U9 §3).

    Unlike figures.duplicate_figure (which requires WRITE on the source figure
    to duplicate at all), only READ access to the source canvas is required
    here — get_canvas(write=False) — so any project member who can merely VIEW
    a shared canvas can still take a personal working copy of it.

    project_id only carries over when the caller ALSO has write access to that
    project (mirrors figure duplicate's "you can only land a copy somewhere you
    could yourself edit"); otherwise the copy is personal (project_id=None)
    rather than silently (re-)sharing it into a project the caller can't edit.

    Panels reference the SAME source figures (figure_id/pinned_version_id
    copied as-is) — no figure is duplicated, so the new canvas renders
    identically to the source immediately. Annotation ids are preserved
    verbatim (they are only unique per-canvas, §U8) but deep-copied so
    mutating one canvas's annotations list can never alias the other's.
    """
    from app.projects import service as project_service

    src = get_canvas(db, canvas_id, owner_id)  # read access is sufficient

    new_project_id = None
    if src.project_id is not None and project_service.can_write_project(db, src.project_id, owner_id):
        new_project_id = src.project_id

    copy_name = (src.name or "Canvas")[: 255 - len(" (copy)")] + " (copy)"

    new_canvas_id = uuid.uuid4()
    canvas = Canvas(
        id=new_canvas_id,
        owner_id=owner_id,
        project_id=new_project_id,
        name=copy_name,
        description=src.description,
        width_mm=src.width_mm,
        height_mm=src.height_mm,
        preset=src.preset,
        background=src.background,
        annotations=copy.deepcopy(src.annotations) if src.annotations else [],
        annotations_rev=0,
    )
    db.add(canvas)

    for p in src.panels:
        db.add(CanvasPanel(
            canvas_id=new_canvas_id,
            figure_id=p.figure_id,
            image_key=p.image_key,
            image_native_width_mm=p.image_native_width_mm,
            image_native_height_mm=p.image_native_height_mm,
            pinned_version_id=p.pinned_version_id,
            x_mm=p.x_mm,
            y_mm=p.y_mm,
            width_mm=p.width_mm,
            height_mm=p.height_mm,
            z_order=p.z_order,
            label=p.label,
            label_visible=p.label_visible,
        ))
    db.commit()
    return canvas_detail(db, new_canvas_id, owner_id)


# ---------------------------------------------------------------- annotations
# U8 — canvas text/shape annotation objects. Structural violations (bad type,
# a required-by-type field missing, a malformed id, or >200 items) fail fast
# with BAD_ANNOTATIONS; numeric ranges are CLAMPED rather than rejected; a
# malformed hex string IS rejected (format error, not an out-of-range number);
# unknown keys are dropped silently by construction (only known keys are ever
# copied into the sanitized entry). Font sizes are absolute pt everywhere
# (design §5) — converted to mm at export time via the existing _PT_TO_MM.
_ANNOTATION_TYPES = {"text", "arrow", "line", "rect", "ellipse"}
_ANNOTATION_MAX_ITEMS = 200
_ANNOTATION_ID_MAX = 64
_ANNOTATION_TEXT_MAX = 500
_ANNOTATION_MM_MIN = -1000.0
_ANNOTATION_MM_MAX = 3000.0
_ANNOTATION_WH_MIN = 0.5
_ANNOTATION_WH_MAX = 2000.0
_ANNOTATION_FONT_MIN = 4.0
_ANNOTATION_FONT_MAX = 72.0
_ANNOTATION_STROKE_MIN = 0.25
_ANNOTATION_STROKE_MAX = 10.0
_ANNOTATION_ALIGN = {"left", "center", "right"}
_ANNOTATION_HEX_RE = re.compile(r"^#[0-9a-fA-F]{6}$")
# XML 1.0 valid char set (minus surrogates): anything outside it (NUL, other
# C0/C1 controls, U+FFFE/U+FFFF, lone surrogates) is stripped from text/id —
# a NUL would 500 the jsonb write (PostgreSQL rejects \\u0000) and any other
# XML-invalid char permanently breaks _compose_canvas_svg exports (libxml2
# refuses the composite SVG).
_ANNOTATION_XML_INVALID_RE = re.compile(
    "[^\x09\x0a\x0d\x20-퟿-�\U00010000-\U0010ffff]"
)


def _ann_float(value) -> float | None:
    """Best-effort numeric coercion; None for anything non-numeric/non-finite
    (bool is deliberately excluded — True/False are not annotation numbers)."""
    if isinstance(value, bool) or not isinstance(value, (int, float, str)):
        return None
    try:
        num = float(value)
    except (TypeError, ValueError):
        return None
    return num if math.isfinite(num) else None


def _ann_hex(value, *, field: str) -> str:
    """Validate a required '#rrggbb' hex string; malformed input is a FORMAT
    error (rejected, not clamped) per the U8 validation contract. fullmatch,
    not match: a $-anchored match still accepts a trailing newline."""
    if not isinstance(value, str) or not _ANNOTATION_HEX_RE.fullmatch(value):
        raise BadRequestError(f"annotation {field} must be a '#rrggbb' hex string", error_code="BAD_ANNOTATIONS")
    return value


def _sanitize_annotations(items: list, canvas_w_mm: float, canvas_h_mm: float) -> list[dict]:
    """Validate + sanitize the ``canvas.annotations`` list (U8 contract).

    ``canvas_w_mm``/``canvas_h_mm`` are accepted for signature parity with the
    design contract but intentionally do NOT clamp coordinates to the sheet —
    an annotation (e.g. a callout arrow) may legitimately extend past the
    sheet edge, same as panel placement already allows (see add_panel's
    x_mm/y_mm comment: "envelope, not canvas-size"). The wide [-1000, 3000] mm
    envelope is the only bound enforced.
    """
    if not isinstance(items, list):
        raise BadRequestError("annotations must be a list", error_code="BAD_ANNOTATIONS")
    if len(items) > _ANNOTATION_MAX_ITEMS:
        raise BadRequestError(
            f"annotations exceeds max of {_ANNOTATION_MAX_ITEMS} items", error_code="BAD_ANNOTATIONS"
        )

    out: list[dict] = []
    seen_ids: set[str] = set()
    for raw in items:
        if not isinstance(raw, dict):
            raise BadRequestError("each annotation must be an object", error_code="BAD_ANNOTATIONS")

        ann_type = raw.get("type")
        if ann_type not in _ANNOTATION_TYPES:
            raise BadRequestError(f"invalid annotation type: {ann_type!r}", error_code="BAD_ANNOTATIONS")

        ann_id = raw.get("id")
        if not isinstance(ann_id, str):
            raise BadRequestError(
                "annotation id must be a non-empty string <=64 chars", error_code="BAD_ANNOTATIONS"
            )
        ann_id = _ANNOTATION_XML_INVALID_RE.sub("", ann_id)
        if not ann_id or len(ann_id) > _ANNOTATION_ID_MAX:
            raise BadRequestError(
                "annotation id must be a non-empty string <=64 chars", error_code="BAD_ANNOTATIONS"
            )
        if ann_id in seen_ids:
            # Duplicate ids corrupt every id-keyed consumer (selection, drag,
            # inspector, undo) for all editors of this canvas — structural
            # violation, fail fast.
            raise BadRequestError(f"duplicate annotation id: {ann_id!r}", error_code="BAD_ANNOTATIONS")
        seen_ids.add(ann_id)

        entry: dict = {"id": ann_id, "type": ann_type}

        try:
            entry["z"] = int(raw.get("z", 0))
        except (TypeError, ValueError):
            entry["z"] = 0

        if ann_type in ("text", "rect", "ellipse"):
            x_mm = _ann_float(raw.get("x_mm"))
            y_mm = _ann_float(raw.get("y_mm"))
            if x_mm is None or y_mm is None:
                raise BadRequestError(
                    f"{ann_type} annotation requires numeric x_mm/y_mm", error_code="BAD_ANNOTATIONS"
                )
            entry["x_mm"] = _clamp(x_mm, _ANNOTATION_MM_MIN, _ANNOTATION_MM_MAX)
            entry["y_mm"] = _clamp(y_mm, _ANNOTATION_MM_MIN, _ANNOTATION_MM_MAX)

        if ann_type in ("rect", "ellipse"):
            w_mm = _ann_float(raw.get("w_mm"))
            h_mm = _ann_float(raw.get("h_mm"))
            if w_mm is None or h_mm is None:
                raise BadRequestError(
                    f"{ann_type} annotation requires numeric w_mm/h_mm", error_code="BAD_ANNOTATIONS"
                )
            entry["w_mm"] = _clamp(w_mm, _ANNOTATION_WH_MIN, _ANNOTATION_WH_MAX)
            entry["h_mm"] = _clamp(h_mm, _ANNOTATION_WH_MIN, _ANNOTATION_WH_MAX)
        elif ann_type == "text":
            # Optional for text: absence means auto-width in the editor.
            if raw.get("w_mm") is not None:
                w_mm = _ann_float(raw.get("w_mm"))
                if w_mm is not None:
                    entry["w_mm"] = _clamp(w_mm, _ANNOTATION_WH_MIN, _ANNOTATION_WH_MAX)
            if raw.get("h_mm") is not None:
                h_mm = _ann_float(raw.get("h_mm"))
                if h_mm is not None:
                    entry["h_mm"] = _clamp(h_mm, _ANNOTATION_WH_MIN, _ANNOTATION_WH_MAX)

        if ann_type in ("arrow", "line"):
            pts = raw.get("points_mm")
            if not isinstance(pts, list) or len(pts) != 4:
                raise BadRequestError(
                    f"{ann_type} annotation requires points_mm [x1,y1,x2,y2]", error_code="BAD_ANNOTATIONS"
                )
            coerced: list[float] = []
            for p in pts:
                pv = _ann_float(p)
                if pv is None:
                    raise BadRequestError(
                        f"{ann_type} annotation points_mm must be numeric", error_code="BAD_ANNOTATIONS"
                    )
                coerced.append(_clamp(pv, _ANNOTATION_MM_MIN, _ANNOTATION_MM_MAX))
            entry["points_mm"] = coerced

        if ann_type == "text":
            text_val = raw.get("text")
            if not isinstance(text_val, str):
                raise BadRequestError("text annotation requires a 'text' string", error_code="BAD_ANNOTATIONS")
            single_line = _ANNOTATION_XML_INVALID_RE.sub("", " ".join(text_val.splitlines())).strip()
            if not single_line:
                raise BadRequestError("text annotation requires non-empty text", error_code="BAD_ANNOTATIONS")
            entry["text"] = single_line[:_ANNOTATION_TEXT_MAX]

            font_pt = _ann_float(raw.get("font_pt")) if raw.get("font_pt") is not None else None
            entry["font_pt"] = (
                _clamp(font_pt, _ANNOTATION_FONT_MIN, _ANNOTATION_FONT_MAX) if font_pt is not None else 10.0
            )

            align = raw.get("align")
            entry["align"] = align if align in _ANNOTATION_ALIGN else "left"

            fill_hex = raw.get("fill_hex")
            entry["fill_hex"] = _ann_hex(fill_hex, field="fill_hex") if fill_hex is not None else "#000000"
        else:
            stroke_hex = raw.get("stroke_hex")
            entry["stroke_hex"] = _ann_hex(stroke_hex, field="stroke_hex") if stroke_hex is not None else "#000000"

            stroke_pt = _ann_float(raw.get("stroke_pt")) if raw.get("stroke_pt") is not None else None
            entry["stroke_pt"] = (
                _clamp(stroke_pt, _ANNOTATION_STROKE_MIN, _ANNOTATION_STROKE_MAX)
                if stroke_pt is not None else 1.0
            )

            if ann_type in ("rect", "ellipse"):
                fill_hex = raw.get("fill_hex")
                entry["fill_hex"] = _ann_hex(fill_hex, field="fill_hex") if fill_hex is not None else None

        out.append(entry)

    return out


# ---------------------------------------------------------------- panels
def _validate_pin(fig: Figure, pinned_version_id: uuid.UUID | None) -> None:
    if pinned_version_id is None:
        return
    if not any(v.id == pinned_version_id for v in fig.versions):
        raise BadRequestError(
            "pinned_version_id does not belong to this figure",
            error_code="BAD_PINNED_VERSION",
        )


def _validate_image_key(image_key: str) -> tuple[float, float]:
    """Validate a client-supplied import key (undo-recreate/duplication path)
    and return the blob's native (w_mm, h_mm).

    The strict key shape plus the existence check bound this to blobs that
    passed the upload pipeline (every file under canvases/imports/ was
    validated + sanitized there) — a key can never point outside that prefix.
    Native size is re-derived from the stored bytes because the key alone
    doesn't carry it.
    """
    if not _IMAGE_KEY_RE.fullmatch(image_key or ""):
        raise BadRequestError("invalid image_key", error_code="BAD_IMAGE_KEY")
    ref = _image_ref(image_key)
    if not storage.exists(ref):
        raise BadRequestError("image_key does not exist", error_code="BAD_IMAGE_KEY")
    blob = storage.read_bytes(ref)
    if _image_ext(image_key) == "svg":
        _sanitized, nw, nh = canvas_images.sanitize_svg(blob)
    else:
        _bytes, nw, nh = canvas_images.validate_raster(blob, _image_ext(image_key))
    return nw, nh


def _fit_image_panel_mm(nw_mm: float, nh_mm: float, canvas_w: float, canvas_h: float) -> tuple[float, float]:
    """Initial panel size for an imported image: native size, shrunk uniformly
    (aspect preserved) if it exceeds ~90% of the sheet, raised uniformly so
    both sides stay >= PANEL_MM_MIN. Mirrors the editor's fitToCanvasMm for
    figure panels. The final per-side clamp only bites on pathological aspect
    ratios (where min-side and max-side constraints conflict)."""
    s = min(1.0, (canvas_w * 0.9) / nw_mm, (canvas_h * 0.9) / nh_mm)
    s = max(s, _PANEL_MM_MIN / nw_mm, _PANEL_MM_MIN / nh_mm)
    return (
        _clamp(nw_mm * s, _PANEL_MM_MIN, _PANEL_MM_MAX),
        _clamp(nh_mm * s, _PANEL_MM_MIN, _PANEL_MM_MAX),
    )


def add_image_panel(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID,
                    content: bytes,
                    x_mm: float | None = None, y_mm: float | None = None,
                    label: str | None = None) -> dict:
    """Upload an external image (SVG/PNG/JPEG) and place it as a new panel.

    Type is sniffed from magic bytes (never the filename); SVG is sanitized
    and the SANITIZED serialization is what gets stored; rasters are decoded
    under the 40M-pixel budget (JPEG re-encoded upright, EXIF dropped).
    ``x_mm``/``y_mm`` are the desired panel CENTER (the client doesn't know
    the final size before the upload); omitted → sheet center. Initial size is
    the image's native physical size fitted to the sheet, aspect preserved.
    """
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    if not content:
        raise BadRequestError("empty image upload", error_code="BAD_IMAGE")

    kind = canvas_images.sniff_kind(content)
    if kind == "svg":
        stored, native_w, native_h = canvas_images.sanitize_svg(content)
    else:
        stored, native_w, native_h = canvas_images.validate_raster(content, kind)

    image_key = f"canvases/imports/{uuid.uuid4().hex}.{kind}"
    media = _IMAGE_MEDIA_TYPES[kind]
    if storage.object_storage_enabled():
        storage.put_bytes(storage.object_key("figures", image_key), stored, content_type=media)
    else:
        path = os.path.join(settings.figures_dir, image_key)
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "wb") as fh:
            fh.write(stored)

    width_mm, height_mm = _fit_image_panel_mm(native_w, native_h, canvas.width_mm, canvas.height_mm)
    center_x = x_mm if x_mm is not None else canvas.width_mm / 2.0
    center_y = y_mm if y_mm is not None else canvas.height_mm / 2.0

    panel = CanvasPanel(
        canvas_id=canvas.id,
        figure_id=None,
        image_key=image_key,
        image_native_width_mm=native_w,
        image_native_height_mm=native_h,
        x_mm=_clamp(center_x - width_mm / 2.0, _PANEL_POS_MIN, _PANEL_POS_MAX),
        y_mm=_clamp(center_y - height_mm / 2.0, _PANEL_POS_MIN, _PANEL_POS_MAX),
        width_mm=width_mm,
        height_mm=height_mm,
        z_order=max((p.z_order for p in canvas.panels), default=-1) + 1,
        label=(label or "").strip()[:8] or None,
    )
    db.add(panel)
    db.commit()
    db.refresh(panel)
    return _panel_response(panel, {}, {})


def add_panel(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID, data: PanelCreate) -> dict:
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    # Exactly one content reference: a figure OR an already-uploaded import
    # blob (the image_key form exists for undo-recreate/duplication — fresh
    # uploads go through add_image_panel).
    if (data.figure_id is None) == (data.image_key is None):
        raise BadRequestError(
            "exactly one of figure_id or image_key is required",
            error_code="BAD_PANEL_CONTENT",
        )

    image_native: tuple[float, float] | None = None
    if data.image_key is not None:
        if data.pinned_version_id is not None:
            raise BadRequestError(
                "image panels cannot pin a figure version", error_code="BAD_PANEL_CONTENT"
            )
        image_native = _validate_image_key(data.image_key)
        fig = None
    else:
        # The figure must be accessible to the caller (owner OR project access);
        # get_figure raises NotFoundError (404) otherwise.
        fig = figures_service.get_figure(db, data.figure_id, owner_id)
        _validate_pin(fig, data.pinned_version_id)

    z_order = data.z_order
    if z_order is None:
        z_order = max((p.z_order for p in canvas.panels), default=-1) + 1

    panel = CanvasPanel(
        canvas_id=canvas.id,
        figure_id=fig.id if fig is not None else None,
        image_key=data.image_key,
        image_native_width_mm=image_native[0] if image_native else None,
        image_native_height_mm=image_native[1] if image_native else None,
        pinned_version_id=data.pinned_version_id,
        # Off-sheet placement envelope [-500, 1000]mm: a panel MAY be parked off
        # the A4 sheet (users place a figure beside the sheet and pull it back).
        # Bounded so it can't be flung unrecoverably far. MUST match the
        # frontend OFF_SHEET_MIN/MAX_MM (CanvasEditor.tsx) so a dragged position
        # never re-snaps server-side.
        x_mm=_clamp(data.x_mm, _PANEL_POS_MIN, _PANEL_POS_MAX),
        y_mm=_clamp(data.y_mm, _PANEL_POS_MIN, _PANEL_POS_MAX),
        width_mm=_clamp(data.width_mm, _PANEL_MM_MIN, _PANEL_MM_MAX),
        height_mm=_clamp(data.height_mm, _PANEL_MM_MIN, _PANEL_MM_MAX),
        z_order=z_order,
        label=data.label,
    )
    db.add(panel)
    db.commit()
    db.refresh(panel)
    current_map = _figure_current_versions(db, [panel.figure_id], owner_id)
    native_map = _version_native_sizes(db, [panel.pinned_version_id or current_map.get(panel.figure_id)])
    return _panel_response(panel, current_map, native_map)


def update_panel(db: Session, canvas_id: uuid.UUID, panel_id: uuid.UUID,
                 owner_id: uuid.UUID, data: dict) -> dict:
    """Mutate placement/size/label/z_order/pin of a panel row ONLY.

    CRITICAL (design §1 "resize = canvas-owned"): this never creates or touches a
    FigureVersion and never invokes the R engine — it writes only to the
    canvas_panels row. Panel width/height changes are geometry the canvas owns.
    """
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    panel = _get_panel(canvas, panel_id)

    if data.get("x_mm") is not None:
        panel.x_mm = _clamp(data["x_mm"], _PANEL_POS_MIN, _PANEL_POS_MAX)
    if data.get("y_mm") is not None:
        panel.y_mm = _clamp(data["y_mm"], _PANEL_POS_MIN, _PANEL_POS_MAX)
    if data.get("width_mm") is not None:
        panel.width_mm = _clamp(data["width_mm"], _PANEL_MM_MIN, _PANEL_MM_MAX)
    if data.get("height_mm") is not None:
        panel.height_mm = _clamp(data["height_mm"], _PANEL_MM_MIN, _PANEL_MM_MAX)
    if data.get("z_order") is not None:
        panel.z_order = data["z_order"]
    if "label" in data:
        panel.label = data["label"]
    if data.get("label_visible") is not None:
        panel.label_visible = data["label_visible"]
    if "pinned_version_id" in data:
        pin = data["pinned_version_id"]
        if pin is not None:
            if panel.figure_id is None:
                raise BadRequestError(
                    "image panels cannot pin a figure version", error_code="BAD_PANEL_CONTENT"
                )
            fig = figures_service.get_figure(db, panel.figure_id, owner_id)
            _validate_pin(fig, pin)
        panel.pinned_version_id = pin

    db.commit()
    db.refresh(panel)
    current_map = _figure_current_versions(db, [panel.figure_id], owner_id)
    native_map = _version_native_sizes(db, [panel.pinned_version_id or current_map.get(panel.figure_id)])
    return _panel_response(panel, current_map, native_map)


def remove_panel(db: Session, canvas_id: uuid.UUID, panel_id: uuid.UUID, owner_id: uuid.UUID) -> None:
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    panel = _get_panel(canvas, panel_id)
    db.delete(panel)
    db.commit()


# ---------------------------------------------------------------- presets
def list_canvas_presets() -> list[dict]:
    """Physical canvas-size presets: ISO paper first, then JOURNAL_SPECS sizes.

    The create dialog seeds its form from the FIRST entry, so list order is the
    default-canvas-size policy: A4 portrait leads. Journal presets follow: for
    each spec, a single-column and double-column canvas (width = column inches
    x 25.4 mm; height = width x 0.72), clamped to [20, 500] mm. No rendering.
    """
    out: list[dict] = [
        # No dimensions in the label — the create dialog appends "(W × H mm)".
        {"key": "a4_portrait", "label": "A4 portrait", "width_mm": 210.0, "height_mm": 297.0},
        {"key": "a4_landscape", "label": "A4 landscape", "width_mm": 297.0, "height_mm": 210.0},
    ]
    for key, spec in JOURNAL_SPECS.items():
        journal = spec.get("journal", key)
        for variant, in_key in (("single", "single_col_in"), ("double", "double_col_in")):
            width_in = spec.get(in_key)
            if not width_in:
                continue
            width_mm = _clamp(width_in * 25.4, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
            height_mm = _clamp(width_mm * 0.72, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
            out.append({
                "key": f"{key}_{variant}",
                "label": f"{journal} — {variant} column ({round(width_mm)} mm)",
                "width_mm": round(width_mm, 2),
                "height_mm": round(height_mm, 2),
            })
    return out


# ============================================================================
# M4 — canvas export (vector composition) + canvas-wide bulk style apply.
#
# CRITICAL (design §1): vector composition ONLY. Each panel is rendered at its
# exact physical mm size by the SAME renderer the preview uses (fonts stay pt),
# then its vector SVG is *nested* inside a parent SVG sized to the canvas. There
# is NO rasterization — the banned rasterGrob/bitmap-stretch path is never used.
# PDF is produced from that composite by rsvg-convert (librsvg), a vector
# SVG→PDF converter that keeps text as text.
# ============================================================================

_EXPORT_PARTS = ("figures", "canvases", "export")

# Panel label (A/B/C…) typography. Absolute pt → mm: the parent SVG's user unit
# is 1 mm (viewBox in mm, width/height in mm), so a physical pt size maps to
# `pt * 25.4 / 72` mm regardless of any panel scaling (§5 font invariance).
_LABEL_PT = 12.0
_PT_TO_MM = 25.4 / 72.0
_LABEL_FONT_MM = round(_LABEL_PT * _PT_TO_MM, 4)
_LABEL_INSET_MM = 1.0  # nudge in from the panel's top-left corner


def _num(value: float) -> str:
    """Format a float for SVG attributes (trim noise, never scientific)."""
    return f"{float(value):.4f}".rstrip("0").rstrip(".") or "0"


def _split_svg(svg_text: str) -> tuple[str, str]:
    """Return (opening-tag attribute string, inner markup) of an SVG document.

    Strips the XML/doctype preamble and the outer <svg> wrapper, keeping the
    children verbatim so vector fidelity (paths, text, defs) is preserved.
    """
    m = re.search(r"<svg\b([^>]*)>(.*)</svg\s*>", svg_text, re.DOTALL | re.IGNORECASE)
    if not m:
        raise BadRequestError("Panel SVG is malformed", error_code="BAD_PANEL_SVG")
    return m.group(1), m.group(2)


def _svg_native_size(attrs: str) -> tuple[float, float]:
    """Native (px) width/height of a panel SVG from its viewBox (preferred) or
    width/height attributes. Used as the nested <svg> viewBox so the panel scales
    cleanly into its mm box."""
    # svglite emits attributes with SINGLE quotes; accept either quote style.
    vb = re.search(r"""viewBox\s*=\s*["']([^"']+)["']""", attrs)
    if vb:
        parts = re.split(r"[\s,]+", vb.group(1).strip())
        if len(parts) == 4:
            try:
                w, h = float(parts[2]), float(parts[3])
                if w > 0 and h > 0:
                    return w, h
            except ValueError:
                pass

    def _dim(name: str) -> float | None:
        mm = re.search(rf"""\b{name}\s*=\s*["']([0-9.]+)""", attrs)
        return float(mm.group(1)) if mm else None

    return (_dim("width") or 100.0, _dim("height") or 100.0)


def _prefix_svg_ids(inner: str, prefix: str) -> str:
    """Scope every id + fragment reference in a panel's inner SVG with `prefix`.

    Nesting multiple independently-rendered SVGs into one document would collide
    on shared ids (svglite emits generic ids for <clipPath>/<defs>/gradients).
    Prefixing definitions AND their references (`url(#id)`, `href="#id"`) with a
    per-panel token keeps each panel's clip-paths/defs self-consistent.
    """
    # svglite uses single-quoted attributes; match either quote style.
    inner = re.sub(r"""\bid=["']([^"']+)["']""", lambda m: f'id="{prefix}{m.group(1)}"', inner)
    inner = re.sub(r"url\(#([^)]+)\)", lambda m: f"url(#{prefix}{m.group(1)})", inner)
    inner = re.sub(
        r"""\b(xlink:href|href)=["']#([^"']+)["']""",
        lambda m: f'{m.group(1)}="#{prefix}{m.group(2)}"',
        inner,
    )
    return inner


_ANNOTATION_ALIGN_ANCHOR = {"left": "start", "center": "middle", "right": "end"}


def _annotation_svg(ann: dict) -> str | None:
    """Render one annotation as an SVG vector primitive in the canvas's mm
    coordinate space (export parity, U8). Returns None for anything malformed
    instead of raising — a stray/legacy annotation must never fail an export
    (§ export parity mandate: "skip, do not crash on malformed legacy items").

    Deliberately re-derives defaults defensively (rather than trusting that
    every stored item passed through ``_sanitize_annotations``) since rows can
    predate a schema tightening or be hand-edited.
    """
    if not isinstance(ann, dict):
        return None
    ann_type = ann.get("type")
    try:
        if ann_type == "text":
            x_mm = float(ann["x_mm"])
            y_mm = float(ann["y_mm"])
            text_val = ann.get("text")
            if not isinstance(text_val, str) or not text_val:
                return None
            font_pt = float(ann.get("font_pt") or 10.0)
            font_mm = font_pt * _PT_TO_MM
            fill_hex = ann.get("fill_hex") or "#000000"
            align = ann.get("align") if ann.get("align") in _ANNOTATION_ALIGN_ANCHOR else "left"
            anchor = "start"
            tx = x_mm
            w_mm = ann.get("w_mm")
            if align != "left" and isinstance(w_mm, (int, float)):
                anchor = _ANNOTATION_ALIGN_ANCHOR[align]
                tx = x_mm + (float(w_mm) / 2.0 if align == "center" else float(w_mm))
            ty = y_mm + 0.8 * font_mm
            return (
                f'<text x="{_num(tx)}" y="{_num(ty)}" '
                f'font-family="Helvetica, Arial, sans-serif" font-size="{_num(font_mm)}" '
                f'fill="{_xml_escape(fill_hex)}" text-anchor="{anchor}">{_xml_escape(text_val)}</text>'
            )

        if ann_type in ("line", "arrow"):
            pts = ann.get("points_mm")
            if not isinstance(pts, list) or len(pts) != 4:
                return None
            x1, y1, x2, y2 = (float(p) for p in pts)
            stroke_hex = ann.get("stroke_hex") or "#000000"
            stroke_pt = float(ann.get("stroke_pt") or 1.0)
            stroke_mm = stroke_pt * _PT_TO_MM
            markup = (
                f'<line x1="{_num(x1)}" y1="{_num(y1)}" x2="{_num(x2)}" y2="{_num(y2)}" '
                f'stroke="{_xml_escape(stroke_hex)}" stroke-width="{_num(stroke_mm)}" stroke-linecap="round"/>'
            )
            if ann_type == "arrow":
                dx, dy = x2 - x1, y2 - y1
                length = math.hypot(dx, dy)
                if length > 1e-6:
                    # No <marker> defs (design §U8): id collisions with the
                    # prefixed panel <defs> and rsvg marker-rendering quirks.
                    head_len = max(2.5 * stroke_mm, 2.0)
                    half_w = 0.6 * head_len
                    ux, uy = dx / length, dy / length
                    px, py = -uy, ux  # unit perpendicular
                    back_x, back_y = x2 - ux * head_len, y2 - uy * head_len
                    p1x, p1y = back_x + px * half_w, back_y + py * half_w
                    p2x, p2y = back_x - px * half_w, back_y - py * half_w
                    markup += (
                        f'<polygon points="{_num(x2)},{_num(y2)} {_num(p1x)},{_num(p1y)} '
                        f'{_num(p2x)},{_num(p2y)}" fill="{_xml_escape(stroke_hex)}"/>'
                    )
            return markup

        if ann_type in ("rect", "ellipse"):
            x_mm = float(ann["x_mm"])
            y_mm = float(ann["y_mm"])
            w_mm = float(ann["w_mm"])
            h_mm = float(ann["h_mm"])
            stroke_hex = ann.get("stroke_hex") or "#000000"
            stroke_pt = float(ann.get("stroke_pt") or 1.0)
            stroke_mm = stroke_pt * _PT_TO_MM
            fill_hex = ann.get("fill_hex")
            fill = _xml_escape(fill_hex) if isinstance(fill_hex, str) and fill_hex else "none"
            if ann_type == "rect":
                return (
                    f'<rect x="{_num(x_mm)}" y="{_num(y_mm)}" width="{_num(w_mm)}" height="{_num(h_mm)}" '
                    f'stroke="{_xml_escape(stroke_hex)}" stroke-width="{_num(stroke_mm)}" fill="{fill}"/>'
                )
            cx, cy = x_mm + w_mm / 2.0, y_mm + h_mm / 2.0
            rx, ry = w_mm / 2.0, h_mm / 2.0
            return (
                f'<ellipse cx="{_num(cx)}" cy="{_num(cy)}" rx="{_num(rx)}" ry="{_num(ry)}" '
                f'stroke="{_xml_escape(stroke_hex)}" stroke-width="{_num(stroke_mm)}" fill="{fill}"/>'
            )
    except (KeyError, TypeError, ValueError):
        return None
    return None


def _annotation_bbox_mm(ann: dict) -> tuple[float, float, float, float] | None:
    """Best-effort mm bounding box (x0, y0, x1, y1) of one annotation, for the
    crop-to-content export. Returns None for malformed items (they are skipped,
    same fail-open policy as _annotation_svg)."""
    if not isinstance(ann, dict):
        return None
    t = ann.get("type")
    try:
        if t in ("rect", "ellipse"):
            x = float(ann["x_mm"]); y = float(ann["y_mm"])
            w = float(ann.get("w_mm") or 0.0); h = float(ann.get("h_mm") or 0.0)
            return (x, y, x + w, y + h)
        if t == "text":
            x = float(ann["x_mm"]); y = float(ann["y_mm"])
            font_mm = float(ann.get("font_pt") or 10.0) * _PT_TO_MM
            # Use the stored width if present; else estimate generously so the
            # crop never clips text (over-inclusion is the safe direction).
            if isinstance(ann.get("w_mm"), (int, float)):
                w = float(ann["w_mm"])
            else:
                w = font_mm * 0.62 * max(1, len(str(ann.get("text") or "")))
            h = float(ann["h_mm"]) if isinstance(ann.get("h_mm"), (int, float)) else font_mm * 1.4
            return (x, y, x + w, y + h)
        if t in ("line", "arrow"):
            pts = ann.get("points_mm")
            if isinstance(pts, list) and len(pts) == 4:
                x1, y1, x2, y2 = (float(p) for p in pts)
                return (min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2))
    except (TypeError, ValueError, KeyError):
        return None
    return None


def _content_bbox_mm(canvas: Canvas) -> tuple[float, float, float, float] | None:
    """Union mm bbox (x0, y0, x1, y1) of every panel + annotation, or None when
    the canvas is empty. Panels use the SAME size clamp as the composite so the
    crop matches what is actually painted."""
    x0s: list[float] = []; y0s: list[float] = []; x1s: list[float] = []; y1s: list[float] = []
    for p in canvas.panels:
        pw = _clamp(p.width_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
        ph = _clamp(p.height_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
        x0s.append(p.x_mm); y0s.append(p.y_mm); x1s.append(p.x_mm + pw); y1s.append(p.y_mm + ph)
    for ann in (canvas.annotations or []):
        bb = _annotation_bbox_mm(ann)
        if bb:
            x0s.append(bb[0]); y0s.append(bb[1]); x1s.append(bb[2]); y1s.append(bb[3])
    if not x0s:
        return None
    return (min(x0s), min(y0s), max(x1s), max(y1s))


def _compose_canvas_svg(db: Session, owner_id: uuid.UUID, canvas: Canvas, crop: bool = False) -> tuple[str, dict[str, str]]:
    """Build the composite canvas SVG by nesting each panel's vector render.

    When ``crop`` is set, the outer <svg> viewBox/size is the tight bounding box
    of all content (panels + annotations) instead of the full A4 sheet — i.e.
    the surrounding margins are removed. Child coordinates are unchanged; the
    viewBox window (min-x/min-y origin) does the cropping, which rsvg-convert
    honours for every format.

    Returns (svg_text, snapshot) where snapshot is {panel_id: version_id} of the
    versions actually rendered (design §5 reproducibility).
    """
    # Full-sheet frame by default; the content bbox when cropping (falls back to
    # the full sheet if the canvas is empty).
    ox, oy = 0.0, 0.0
    w_mm = _clamp(canvas.width_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    h_mm = _clamp(canvas.height_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    if crop:
        bbox = _content_bbox_mm(canvas)
        if bbox is not None:
            bx0, by0, bx1, by1 = bbox
            ox, oy = bx0, by0
            w_mm = max(1.0, bx1 - bx0)
            h_mm = max(1.0, by1 - by0)

    current_map = _figure_current_versions(db, [p.figure_id for p in canvas.panels], owner_id)
    # Paint order: ascending z_order (ties by id) → later panels drawn on top (§2).
    panels = sorted(canvas.panels, key=lambda p: (p.z_order, str(p.id)))

    parts: list[str] = [
        '<?xml version="1.0" encoding="UTF-8"?>',
        (
            '<svg xmlns="http://www.w3.org/2000/svg" '
            'xmlns:xlink="http://www.w3.org/1999/xlink" version="1.1" '
            f'width="{_num(w_mm)}mm" height="{_num(h_mm)}mm" '
            f'viewBox="{_num(ox)} {_num(oy)} {_num(w_mm)} {_num(h_mm)}">'
        ),
    ]
    # Background: opaque white unless the canvas is explicitly transparent (§2).
    # When cropping, the rect covers the bbox (origin ox/oy), not 0,0.
    if canvas.background != "transparent":
        parts.append(f'<rect x="{_num(ox)}" y="{_num(oy)}" width="{_num(w_mm)}" height="{_num(h_mm)}" fill="#ffffff"/>')

    snapshot: dict[str, str] = {}
    for idx, panel in enumerate(panels):
        pw_mm = _clamp(panel.width_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
        ph_mm = _clamp(panel.height_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)

        effective = None
        if panel.image_key:
            # Imported image: SVG imports nest as vector (same path as figure
            # panels — the stored markup is already sanitized); raster imports
            # embed as a data-URI <image>, which rsvg resolves locally with no
            # network access. A missing/corrupt blob skips the panel (same
            # fail-open-per-panel policy as an inaccessible figure).
            try:
                blob = storage.read_bytes(_image_ref(panel.image_key))
            except (OSError, RuntimeError, ValueError):
                continue
            if _image_ext(panel.image_key) == "svg":
                try:
                    attrs, inner = _split_svg(blob.decode("utf-8"))
                except (BadRequestError, UnicodeDecodeError):
                    continue
                native_w, native_h = _svg_native_size(attrs)
                inner = _prefix_svg_ids(inner, f"p{idx}_")
                parts.append(
                    f'<svg x="{_num(panel.x_mm)}" y="{_num(panel.y_mm)}" '
                    f'width="{_num(pw_mm)}" height="{_num(ph_mm)}" '
                    f'viewBox="0 0 {_num(native_w)} {_num(native_h)}" '
                    f'preserveAspectRatio="none" overflow="hidden">{inner}</svg>'
                )
            else:
                media = _IMAGE_MEDIA_TYPES[_image_ext(panel.image_key)]
                b64 = base64.b64encode(blob).decode("ascii")
                parts.append(
                    f'<image x="{_num(panel.x_mm)}" y="{_num(panel.y_mm)}" '
                    f'width="{_num(pw_mm)}" height="{_num(ph_mm)}" '
                    f'preserveAspectRatio="none" '
                    f'xlink:href="data:{media};base64,{b64}"/>'
                )
        else:
            effective = panel.pinned_version_id or current_map.get(panel.figure_id)
            if effective is None:
                continue  # figure has no version yet → nothing to render

            # Reuse the exact physical-size vector renderer (fonts stay pt, §5).
            try:
                ref, _layout, _cached = _render_preview_ref(
                    db,
                    owner_id,
                    PreviewRenderRequest(
                        figure_id=panel.figure_id,
                        version_id=effective,
                        width_mm=pw_mm,
                        height_mm=ph_mm,
                    ),
                )
            except NotFoundError:
                # The exporting user can't access this panel's figure (e.g. the
                # canvas was moved into a project while a panel still references a
                # personal figure). Fail-closed per PANEL — skip it — instead of
                # hard-404ing a collaborator's ENTIRE export on the first miss.
                continue
            svg_text = storage.read_bytes(ref).decode("utf-8")
            attrs, inner = _split_svg(svg_text)
            native_w, native_h = _svg_native_size(attrs)
            inner = _prefix_svg_ids(inner, f"p{idx}_")

            parts.append(
                f'<svg x="{_num(panel.x_mm)}" y="{_num(panel.y_mm)}" '
                f'width="{_num(pw_mm)}" height="{_num(ph_mm)}" '
                f'viewBox="0 0 {_num(native_w)} {_num(native_h)}" '
                f'preserveAspectRatio="none" overflow="hidden">{inner}</svg>'
            )

        if panel.label_visible and panel.label:
            # Absolute-pt bold label pinned to the panel's top-left corner.
            ty = panel.y_mm + _LABEL_INSET_MM + _LABEL_FONT_MM
            tx = panel.x_mm + _LABEL_INSET_MM
            parts.append(
                f'<text x="{_num(tx)}" y="{_num(ty)}" '
                f'font-family="Helvetica, Arial, sans-serif" '
                f'font-size="{_num(_LABEL_FONT_MM)}" font-weight="bold" '
                f'fill="#000000">{_xml_escape(panel.label)}</text>'
            )

        # Snapshot records figure VERSIONS (§5 reproducibility) — image panels
        # have no version and are immutable blobs, so they don't appear.
        if effective is not None:
            snapshot[str(panel.id)] = str(effective)

    # U8: annotations are emitted AFTER every panel — they always paint ABOVE
    # panels (V1) — sorted by (z, id) for stable, deterministic paint order.
    annotations = sorted(
        (a for a in (canvas.annotations or []) if isinstance(a, dict)),
        key=lambda a: (a.get("z", 0) if isinstance(a.get("z"), (int, float)) else 0, str(a.get("id", ""))),
    )
    for ann in annotations:
        markup = _annotation_svg(ann)
        if markup:
            parts.append(markup)

    parts.append("</svg>")
    return "\n".join(parts), snapshot


def _persist_export(canvas_id: uuid.UUID, data: bytes, ext: str, content_type: str) -> str:
    """Write the composed export artifact to storage and return a served URL."""
    filename = f"{canvas_id}_{uuid.uuid4().hex}.{ext}"
    if storage.object_storage_enabled():
        key = storage.object_key(*_EXPORT_PARTS, filename)
        ref = storage.put_bytes(key, data, content_type=content_type)
        return figures_service._url(ref)
    out_path = os.path.join(settings.figures_dir, "canvases", "export", filename)
    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    with open(out_path, "wb") as fh:
        fh.write(data)
    return figures_service._url(out_path)


_RASTER_DPI_CHOICES = {300, 600}
# Raster pixel budget (U9 §2 hardening). Pixel size is ceil(mm/25.4*dpi) per
# side; this caps the rsvg child (~4 B/px RGBA) AND the in-process Pillow TIFF
# re-encode (~8 B/px peak — measured: 139.5M px -> ~1.1 GB RSS, enough to OOM
# the prod box under 2-3 concurrent exports). 40M px keeps the Pillow peak
# ~310 MB while still allowing every legitimate combination: the max 500x500mm
# canvas @300 dpi is 34.9M px and an A4 sheet @600 dpi is 34.8M px.
_RASTER_MAX_PIXELS = 40_000_000


def _svg_to_png_bytes(svg_text: str, dpi: int) -> bytes:
    """Rasterize a physical-unit (mm) composite SVG to PNG at `dpi` via
    rsvg-convert (U9 §2).

    The composite SVG carries width/height in mm (e.g. ``width="210mm"``) with
    a 1mm-per-user-unit viewBox. Passing ``-d/-p <dpi>`` WITHOUT
    ``--width``/``--height`` lets rsvg-convert derive the pixel size itself
    from those physical dimensions rather than us precomputing one — verified
    in-container against the installed rsvg-convert 2.52.5: output pixel size
    is ``ceil(mm / 25.4 * dpi)`` (NOT round-to-nearest), e.g. a 210mm-wide
    sheet -> 2481px at 300dpi (ceil(2480.315...)) and 4961px at 600dpi.
    """
    with tempfile.TemporaryDirectory(prefix="labplot_canvas_raster_") as td:
        svg_path = os.path.join(td, "composite.svg")
        png_path = os.path.join(td, "out.png")
        with open(svg_path, "w", encoding="utf-8") as fh:
            fh.write(svg_text)
        # Hardened subprocess: no shell, fixed literal args + trusted temp
        # paths, bounded runtime (mirrors the PDF branch).
        try:
            subprocess.run(
                ["rsvg-convert", "-f", "png", "-d", str(dpi), "-p", str(dpi), "-o", png_path, svg_path],
                check=True,
                capture_output=True,
                timeout=60,
            )
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
            raise BadRequestError("PNG conversion failed", error_code="PNG_EXPORT_FAILED") from exc
        with open(png_path, "rb") as fh:
            png_bytes = fh.read()
    if not png_bytes:
        raise BadRequestError("PNG conversion produced no output", error_code="PNG_EXPORT_FAILED")
    return png_bytes


def _stamp_png_dpi(png_bytes: bytes, dpi: int) -> bytes:
    """Re-encode a PNG with a pHYs resolution chunk (U9 review F1): rsvg's
    cairo writer emits none, so a '300 dpi' PNG would open as 72 dpi in
    Photoshop/journal submission checkers — the exact use case the dpi option
    exists for. Pillow writes pHYs when `dpi` is passed to save()."""
    from PIL import Image

    Image.MAX_IMAGE_PIXELS = _RASTER_MAX_PIXELS  # export_canvas guards to the same budget
    with Image.open(io.BytesIO(png_bytes)) as im:
        out = io.BytesIO()
        im.save(out, format="PNG", dpi=(dpi, dpi))
        return out.getvalue()


def _png_bytes_to_tiff_bytes(png_bytes: bytes, dpi: int) -> bytes:
    """Convert PNG bytes to LZW-compressed TIFF bytes via Pillow, preserving
    DPI metadata (U9 §2). Reuses the rsvg-produced PNG as the raster source
    rather than a second rsvg invocation. Local import: Pillow is only needed
    on this raster export path."""
    from PIL import Image

    # export_canvas guards inputs to <= _RASTER_MAX_PIXELS; make Pillow enforce
    # the same ceiling (its default only WARNS at ~89.5M px and errors at 2x),
    # so the budget cannot be bypassed by any future caller.
    Image.MAX_IMAGE_PIXELS = _RASTER_MAX_PIXELS
    with Image.open(io.BytesIO(png_bytes)) as im:
        out = io.BytesIO()
        im.save(out, format="TIFF", compression="tiff_lzw", dpi=(dpi, dpi))
        return out.getvalue()


# DPI used to rasterise each panel figure into the PPTX. Slides don't need
# 600dpi; 200 keeps each panel PNG well under the raster budget while staying
# crisp on a projector.
_PPTX_PANEL_DPI = 200


def _build_pptx_bytes(db: Session, owner_id: uuid.UUID, canvas: Canvas, crop: bool = False) -> tuple[bytes, dict[str, str]]:
    """Build a one-slide PPTX where each figure is its OWN movable picture and
    annotations are native PowerPoint shapes, sized so a slide millimetre maps
    to a canvas millimetre. Returns (pptx_bytes, snapshot).

    python-pptx is imported lazily so the module never hard-depends on it — an
    absent package surfaces as a clear 501 on THIS path only.
    """
    try:
        from pptx import Presentation
        from pptx.util import Mm, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.text import PP_ALIGN
    except ImportError as exc:  # pragma: no cover - depends on deploy image
        raise AppError(
            status_code=501,
            detail="PPTX export requires the python-pptx package, which is not installed.",
            error_code="PPTX_EXPORT_UNAVAILABLE",
        ) from exc

    def _rgb(hexstr, default=(0, 0, 0)):
        s = str(hexstr or "").lstrip("#")
        if len(s) == 6:
            try:
                return RGBColor.from_string(s.upper())
            except ValueError:
                pass
        return RGBColor(*default)

    # Slide = full sheet, or the content bbox when cropping. ox/oy shift every
    # placed object into slide space.
    ox, oy = 0.0, 0.0
    slide_w = _clamp(canvas.width_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    slide_h = _clamp(canvas.height_mm, _CANVAS_MM_MIN, _CANVAS_MM_MAX)
    if crop:
        bbox = _content_bbox_mm(canvas)
        if bbox is not None:
            ox, oy = bbox[0], bbox[1]
            slide_w = max(1.0, bbox[2] - bbox[0])
            slide_h = max(1.0, bbox[3] - bbox[1])

    prs = Presentation()
    prs.slide_width = Mm(slide_w)
    prs.slide_height = Mm(slide_h)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    # White page fill via the SLIDE BACKGROUND (a slide-level property), NOT a
    # shape. A background rectangle is selectable/clickable in PowerPoint, so it
    # gets in the way when the user just wants to grab a figure — only the
    # figures/annotations should be objects. slide.background writes <p:bg> into
    # the slide XML and is not part of the shape tree (verified: shape_count 0).
    if canvas.background != "transparent":
        try:
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        except Exception:  # noqa: BLE001 - background is decorative, never fail the export
            pass

    current_map = _figure_current_versions(db, [p.figure_id for p in canvas.panels], owner_id)
    panels = sorted(canvas.panels, key=lambda p: (p.z_order, str(p.id)))
    snapshot: dict[str, str] = {}
    for panel in panels:
        pw = _clamp(panel.width_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
        ph = _clamp(panel.height_mm, _PANEL_MM_MIN, _PANEL_MM_MAX)
        effective = None
        if panel.image_key:
            # Imported image: PNG/JPEG insert natively (python-pptx reads
            # both); SVG is wrapped at the panel's physical mm size and
            # rasterized through the same rsvg path as figure panels.
            try:
                blob = storage.read_bytes(_image_ref(panel.image_key))
            except (OSError, RuntimeError, ValueError):
                continue
            if _image_ext(panel.image_key) == "svg":
                try:
                    attrs, inner = _split_svg(blob.decode("utf-8"))
                except (BadRequestError, UnicodeDecodeError):
                    continue
                native_w, native_h = _svg_native_size(attrs)
                wrapper = (
                    '<svg xmlns="http://www.w3.org/2000/svg" '
                    'xmlns:xlink="http://www.w3.org/1999/xlink" '
                    f'width="{_num(pw)}mm" height="{_num(ph)}mm" '
                    f'viewBox="0 0 {_num(pw)} {_num(ph)}">'
                    f'<svg x="0" y="0" width="{_num(pw)}" height="{_num(ph)}" '
                    f'viewBox="0 0 {_num(native_w)} {_num(native_h)}" '
                    f'preserveAspectRatio="none" overflow="hidden">{inner}</svg></svg>'
                )
                picture_bytes = _svg_to_png_bytes(wrapper, _PPTX_PANEL_DPI)
            else:
                picture_bytes = blob
            try:
                slide.shapes.add_picture(io.BytesIO(picture_bytes),
                                         Mm(panel.x_mm - ox), Mm(panel.y_mm - oy), Mm(pw), Mm(ph))
            except Exception:  # noqa: BLE001 - skip an unreadable image, keep the export
                continue
        else:
            effective = panel.pinned_version_id or current_map.get(panel.figure_id)
            if effective is None:
                continue
            try:
                ref, _layout, _cached = _render_preview_ref(
                    db, owner_id,
                    PreviewRenderRequest(figure_id=panel.figure_id, version_id=effective, width_mm=pw, height_mm=ph),
                )
            except NotFoundError:
                continue  # inaccessible panel figure — skip, don't fail the whole export
            svg_text = storage.read_bytes(ref).decode("utf-8")
            png_bytes = _svg_to_png_bytes(svg_text, _PPTX_PANEL_DPI)
            slide.shapes.add_picture(io.BytesIO(png_bytes), Mm(panel.x_mm - ox), Mm(panel.y_mm - oy), Mm(pw), Mm(ph))
        if panel.label_visible and panel.label:
            try:
                tb = slide.shapes.add_textbox(Mm(panel.x_mm - ox + _LABEL_INSET_MM), Mm(panel.y_mm - oy),
                                              Mm(20), Mm(_LABEL_FONT_MM * 1.8))
                tf = tb.text_frame
                tf.word_wrap = False
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                run = tf.paragraphs[0].add_run()
                run.text = panel.label
                run.font.bold = True
                run.font.size = Pt(_LABEL_FONT_MM / _PT_TO_MM)
                run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception:  # noqa: BLE001 - label is decorative
                pass
        if effective is not None:
            snapshot[str(panel.id)] = str(effective)

    annotations = sorted(
        (a for a in (canvas.annotations or []) if isinstance(a, dict)),
        key=lambda a: (a.get("z", 0) if isinstance(a.get("z"), (int, float)) else 0, str(a.get("id", ""))),
    )
    for ann in annotations:
        t = ann.get("type")
        try:
            if t in ("rect", "ellipse"):
                x = float(ann["x_mm"]); y = float(ann["y_mm"])
                w = float(ann.get("w_mm") or 0.0); h = float(ann.get("h_mm") or 0.0)
                if w <= 0 or h <= 0:
                    continue
                shp = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL if t == "ellipse" else MSO_SHAPE.RECTANGLE,
                    Mm(x - ox), Mm(y - oy), Mm(w), Mm(h),
                )
                shp.shadow.inherit = False
                fill_hex = ann.get("fill_hex")
                if fill_hex and str(fill_hex).lower() not in ("none", "transparent", ""):
                    shp.fill.solid(); shp.fill.fore_color.rgb = _rgb(fill_hex)
                else:
                    shp.fill.background()
                stroke_hex = ann.get("stroke_hex")
                if stroke_hex:
                    shp.line.color.rgb = _rgb(stroke_hex)
                    shp.line.width = Pt(float(ann.get("stroke_pt") or 1.0))
                else:
                    shp.line.fill.background()
            elif t == "text":
                text_val = ann.get("text")
                if not isinstance(text_val, str) or not text_val:
                    continue
                x = float(ann["x_mm"]); y = float(ann["y_mm"])
                font_pt = float(ann.get("font_pt") or 10.0)
                w = float(ann["w_mm"]) if isinstance(ann.get("w_mm"), (int, float)) \
                    else font_pt * _PT_TO_MM * 0.62 * max(1, len(text_val))
                h = float(ann["h_mm"]) if isinstance(ann.get("h_mm"), (int, float)) else font_pt * _PT_TO_MM * 1.8
                tb = slide.shapes.add_textbox(Mm(x - ox), Mm(y - oy), Mm(max(w, 5.0)), Mm(max(h, 5.0)))
                tf = tb.text_frame
                tf.word_wrap = False
                tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                para = tf.paragraphs[0]
                para.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(ann.get("align"), PP_ALIGN.LEFT)
                run = para.add_run()
                run.text = text_val
                run.font.size = Pt(font_pt)
                run.font.color.rgb = _rgb(ann.get("fill_hex"))
                run.font.name = "Helvetica"
            elif t in ("line", "arrow"):
                pts = ann.get("points_mm")
                if not (isinstance(pts, list) and len(pts) == 4):
                    continue
                x1, y1, x2, y2 = (float(p) for p in pts)
                cxn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                                 Mm(x1 - ox), Mm(y1 - oy), Mm(x2 - ox), Mm(y2 - oy))
                cxn.line.color.rgb = _rgb(ann.get("stroke_hex"))
                cxn.line.width = Pt(float(ann.get("stroke_pt") or 1.0))
                if t == "arrow":
                    try:
                        from pptx.oxml.ns import qn
                        ln = cxn.line._get_or_add_ln()
                        ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle"}))
                    except Exception:  # noqa: BLE001 - arrowhead is cosmetic
                        pass
        except (TypeError, ValueError, KeyError):
            continue  # malformed annotation — skip (export-parity fail-open)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue(), snapshot


def export_canvas(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID, fmt: str, dpi: int = 300, crop: bool = False) -> dict:
    """Export the canvas as svg, pdf, png, tiff, or pptx (§1, U9 §2).

    svg/pdf stay pure vector composition (§1: never bitmap-stretch a panel).
    png/tiff RASTERIZE that SAME composite via rsvg-convert at `dpi` — still a
    single controlled conversion of the whole composed sheet, not an ad hoc
    per-panel bitmap stretch. tiff is derived from the rendered png (one rsvg
    invocation, then a Pillow re-encode) rather than a second rsvg call.

    Records `{panel_id: version_id}` into `canvas.export_snapshot` for
    reproducibility (§5) and commits, for every format.
    """
    fmt = (fmt or "svg").lower()
    if fmt not in {"svg", "pdf", "png", "tiff", "pptx"}:
        raise BadRequestError("format must be one of 'svg', 'pdf', 'png', 'tiff', 'pptx'", error_code="BAD_EXPORT_FORMAT")
    if fmt in {"png", "tiff"} and dpi not in _RASTER_DPI_CHOICES:
        raise BadRequestError("dpi must be 300 or 600", error_code="BAD_EXPORT_DPI")

    # Fail fast if a converter-backed format is requested but the vector
    # converter is unavailable — never silently fall back to a lesser path.
    # pptx rasterises each panel through the same rsvg path.
    if fmt in {"pdf", "png", "tiff", "pptx"} and shutil.which("rsvg-convert") is None:
        raise AppError(
            status_code=501,
            detail=f"{fmt.upper()} export requires librsvg (rsvg-convert), which is not installed.",
            error_code="PDF_EXPORT_UNAVAILABLE" if fmt == "pdf" else "RASTER_EXPORT_UNAVAILABLE",
        )

    canvas = get_canvas(db, canvas_id, owner_id)

    # The rasterised size is the CROPPED content bbox when cropping (smaller),
    # else the full sheet — so the pixel-budget guard matches the real output.
    guard_w = float(canvas.width_mm)
    guard_h = float(canvas.height_mm)
    if crop:
        bbox = _content_bbox_mm(canvas)
        if bbox is not None:
            guard_w = max(1.0, bbox[2] - bbox[0])
            guard_h = max(1.0, bbox[3] - bbox[1])

    if fmt in {"png", "tiff"}:
        # Guard BEFORE composing (per-panel renders are the expensive part):
        # ceil matches rsvg-convert's observed mm->px rounding (see the
        # _svg_to_png_bytes docstring).
        px_w = math.ceil(guard_w / 25.4 * dpi)
        px_h = math.ceil(guard_h / 25.4 * dpi)
        if px_w * px_h > _RASTER_MAX_PIXELS:
            raise BadRequestError(
                f"Raster export would be {px_w}x{px_h}px, over the "
                f"{_RASTER_MAX_PIXELS // 1_000_000}M-pixel limit — use 300 dpi or a smaller canvas.",
                error_code="RASTER_TOO_LARGE",
            )

    # pptx composes its OWN document (per-panel pictures), not the flat SVG —
    # handle it early so we don't render every panel twice.
    if fmt == "pptx":
        pptx_bytes, snapshot = _build_pptx_bytes(db, owner_id, canvas, crop=crop)
        url = _persist_export(
            canvas_id, pptx_bytes, "pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        canvas.export_snapshot = snapshot
        db.commit()
        return {"url": url, "format": fmt, "dpi": None, "snapshot": snapshot}

    composite_svg, snapshot = _compose_canvas_svg(db, owner_id, canvas, crop=crop)

    result_dpi: int | None = None
    if fmt == "svg":
        url = _persist_export(canvas_id, composite_svg.encode("utf-8"), "svg", "image/svg+xml")
    elif fmt == "pdf":
        with tempfile.TemporaryDirectory(prefix="labplot_canvas_export_") as td:
            svg_path = os.path.join(td, "composite.svg")
            pdf_path = os.path.join(td, "out.pdf")
            with open(svg_path, "w", encoding="utf-8") as fh:
                fh.write(composite_svg)
            # Hardened subprocess: no shell, fixed literal args + trusted temp
            # paths, bounded runtime.
            try:
                subprocess.run(
                    ["rsvg-convert", "-f", "pdf", "-o", pdf_path, svg_path],
                    check=True,
                    capture_output=True,
                    timeout=60,
                )
            except (subprocess.CalledProcessError, subprocess.TimeoutExpired) as exc:
                raise BadRequestError("PDF conversion failed", error_code="PDF_EXPORT_FAILED") from exc
            with open(pdf_path, "rb") as fh:
                pdf_bytes = fh.read()
        if not pdf_bytes:
            raise BadRequestError("PDF conversion produced no output", error_code="PDF_EXPORT_FAILED")
        url = _persist_export(canvas_id, pdf_bytes, "pdf", "application/pdf")
    elif fmt == "png":
        result_dpi = dpi
        png_bytes = _stamp_png_dpi(_svg_to_png_bytes(composite_svg, dpi), dpi)
        url = _persist_export(canvas_id, png_bytes, "png", "image/png")
    else:  # tiff
        result_dpi = dpi
        png_bytes = _svg_to_png_bytes(composite_svg, dpi)
        tiff_bytes = _png_bytes_to_tiff_bytes(png_bytes, dpi)
        url = _persist_export(canvas_id, tiff_bytes, "tiff", "image/tiff")

    # Snapshot the versions actually composed (§5), then commit.
    canvas.export_snapshot = snapshot
    db.commit()
    return {"url": url, "format": fmt, "dpi": result_dpi, "snapshot": snapshot}


def apply_canvas_style(db: Session, canvas_id: uuid.UUID, owner_id: uuid.UUID,
                       source_figure_id: uuid.UUID) -> dict:
    """Apply the source figure's STYLE-ONLY options to every OTHER panel figure.

    Delegates to `figures_service.bulk_apply_style`, so each target figure gets a
    NEW version (content ⇒ version bump, design decision 3). The source must be
    one of the canvas's own panel figures. Returns {updated, skipped}.
    """
    canvas = get_canvas(db, canvas_id, owner_id, write=True)
    # Image panels carry no figure/style — they neither source nor receive.
    panel_figure_ids = [p.figure_id for p in canvas.panels if p.figure_id is not None]
    if source_figure_id not in panel_figure_ids:
        raise BadRequestError(
            "source_figure_id must be one of the canvas's panel figures",
            error_code="BAD_SOURCE_FIGURE",
        )
    targets = list({fid for fid in panel_figure_ids if fid != source_figure_id})
    result = figures_service.bulk_apply_style(db, source_figure_id, targets, owner_id)
    return {"updated": result.get("updated", []), "skipped": result.get("skipped", [])}
